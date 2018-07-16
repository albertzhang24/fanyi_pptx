import json
import sys, os
from hashlib import md5
import sys, hashlib, random, json, string
import urllib.request, urllib.parse, urllib.error
import http.client

import requests
import boto3
from botocore.client import Config

from pptx import Presentation

__author__ = 'albertzhang'

appid = '20180531000169614'
secretKey = 'IRP4F_z0fxRscAU5aMt_'
salt = random.randint(32768, 65536)


BUCKET_NAME = 'fanyi-pptx'


def home(event, context):
    index_content = ''
    # Get index file and show to user.
    with open('templates/index.html', 'r') as index_file:
        index_content = index_file.read()

    return {
        'statusCode': 200,
        'headers': {
            'Content-Type': 'text/html',
        },
        'body': index_content
    }


def translate(event, context):
    # We are expecting the user to send us an url.
    untranslated_pptx_url = json.loads(event['body'])['url']

    # Download the package and write to tmp location.
    tmp_path = os.path.join('/tmp', os.path.basename(untranslated_pptx_url))
    untranslated_pptx_response = requests.get(untranslated_pptx_url, stream=True)

    with open(tmp_path, 'wb') as tmp_path_file:
        for block in untranslated_pptx_response.iter_content(1024):
            tmp_path_file.write(block)

    # Translate the PPTX and upload to S3 converted bucket.
    translated_filepath = translate_powerpoint(tmp_path)
    translated_url = upload_to_s3(translated_filepath)
    print(translated_url)
    # Return the converted S3 URL.
    return {
        'statusCode': 200,
        'headers': {
            'Content-Type': 'application/json',
        },
        'body': json.dumps({
            'url': translated_url
        })
    }


def translate_powerpoint(input_filename):
    # source_target = request.form['source_target']
    sourceLang = 'en'
    targetLang = 'zh'
    
    # if source_target[:source_target.find("_")] == 'en':
    #     sourceLang = 'en'
    #     targetLang = 'zh'
    # else:
    #     sourceLang = 'zh'
    #     targetLang = 'en'

    source_ppt_str = input_filename
    ppt = Presentation(input_filename)

    target_ppt = source_ppt_str[:source_ppt_str.find(
        '.pptx')] + "_" + targetLang + ".pptx"

    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = _translate(run.text, sourceLang, targetLang)

    ppt.save(target_ppt)

    return target_ppt


def _translate(q, sourceLang, targetLang):
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')

    m1 = hashlib.md5()
    m1.update((appid+q+str(salt)+secretKey).encode('utf-8'))
    sign = m1.hexdigest()
    url='/api/trans/vip/translate?appid=%s&q=%s&from=%s&to=%s&salt=%s&sign=%s' % (
        appid, urllib.parse.quote(q), sourceLang, targetLang, str(salt), sign)

    count = 0

    try:
        httpClient.request('GET', url)
        response = httpClient.getresponse()

        r = str(response.read())[2:-1]
        j = json.loads(r)

        if 'trans_result' in j:
            return j['trans_result'][-1]['dst'].encode('ascii').decode('unicode-escape')
        elif 'error_msg' in j:
            count += 1
            return '***Error (%s:%s): %s' % (j['error_code'], j['error_msg'], q)
        else:
            count += 1
            return '***Unexpected response for %s: %s' %(q, r)

    except Exception as e:
        count += 1
        return '***Exception (%s): %s' % (e,q)


def upload_to_s3(filepath):
    client = boto3.client('cognito-identity')
    credentials = client.get_credentials_for_identity(
        IdentityId='us-west-2:7e280494-3998-488e-832e-14bf8d4abde7',
    )

    s3 = boto3.resource('s3',
        aws_access_key_id=credentials['Credentials']['AccessKeyId'],
        aws_secret_access_key=credentials['Credentials']['SecretKey'],
        aws_session_token=credentials['Credentials']['SessionToken'],
        config=Config(signature_version='s3v4')
    )

    # Upload a new file
    key = os.path.basename(filepath)
    s3.Bucket(BUCKET_NAME).put_object(
        ACL='public-read',
        Key=key,
        Body=open(filepath, 'rb')
    )

    s3_client = boto3.client('s3')
    return s3_client.generate_presigned_url(
        'get_object', Params={
            'Bucket': BUCKET_NAME,
            'Key': key
        }, ExpiresIn=3600
    )
